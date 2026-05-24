import os
import fitz  # PyMuPDF
from pptx import Presentation
from google import genai
import backend.database as db

class RAGManager:
    def __init__(self):
        # We assume GEMINI_API_KEY or GOOGLE_API_KEY is loaded in environment variables
        api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
        if not api_key:
            print("WARNING: GEMINI_API_KEY or GOOGLE_API_KEY environment variable is not set!")
            self.genai_client = None
            return
            
        self.genai_client = genai.Client(api_key=api_key)

    def extract_text_from_pdf(self, file_path: str):
        """Extract pages from PDF with page numbers."""
        pages = []
        try:
            doc = fitz.open(file_path)
            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    pages.append({
                        "text": text.strip(),
                        "page_number": i + 1
                    })
        except Exception as e:
            print(f"Failed to parse PDF {file_path}: {e}")
        return pages

    def extract_text_from_pptx(self, file_path: str):
        """Extract slides from PPTX with slide numbers."""
        slides = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                text = "\n".join(text_parts)
                if text.strip():
                    slides.append({
                        "text": text.strip(),
                        "page_number": i + 1
                    })
        except Exception as e:
            print(f"Failed to parse PPTX {file_path}: {e}")
        return slides

    def ingest_file(self, file_path: str, subject: str = "General"):
        """Extracts text, aggregates it by page/slide, and returns the metadata and full structured text."""
        filename = os.path.basename(file_path)
        ext = os.path.splitext(filename)[1].lower()
        
        # 1. Extract Text
        if ext == ".pdf":
            print(f"Extracting raw text from PDF: {filename}...")
            raw_pages = self.extract_text_from_pdf(file_path)
            
            items = []
            for page in raw_pages:
                raw_text = page["text"]
                p_num = page["page_number"]
                
                formatted_text = raw_text
                if self.genai_client and raw_text.strip():
                    try:
                        print(f"Converting PDF page {p_num} text to Markdown using Gemini...")
                        prompt = (
                            f"You are a document processing assistant. Convert the following raw text extracted from "
                            f"Page {p_num} of a PDF into clean, well-formatted, rich Markdown.\n\n"
                            f"Rules:\n"
                            f"1. Preserve the structural layout of the page: format lists as Markdown lists, tables as Markdown tables, and section headings as headings (e.g. ## Header).\n"
                            f"2. DO NOT summarize, compress, omit, or modify any facts, numbers, dates, formulas, or text content. Retain every detail exactly as it is.\n"
                            f"3. Do not include any meta-commentary, introductory text, or Markdown codeblock fences (like ```markdown). Output ONLY the raw formatted Markdown.\n\n"
                            f"Raw Text:\n{raw_text}"
                        )
                        response = self.genai_client.models.generate_content(
                            model="gemini-2.5-flash",
                            contents=prompt
                        )
                        if response.text:
                            formatted_text = response.text.strip()
                    except Exception as e:
                        print(f"Gemini markdown conversion failed for page {p_num}: {e}. Falling back to raw text.")
                
                items.append({
                    "text": formatted_text,
                    "page_number": p_num
                })
                
            # Formulate full text
            full_content_parts = []
            for item in items:
                full_content_parts.append(f"--- Page {item['page_number']} ---\n{item['text']}\n")
            full_content = "\n".join(full_content_parts)
            
            print(f"Parsed {len(items)} pages from PDF.")
            return len(items), full_content, items
            
        elif ext in [".pptx", ".ppt"]:
            items = self.extract_text_from_pptx(file_path)
            if not items:
                return 0, "", []
                
            full_content_parts = []
            for item in items:
                p_num = item["page_number"]
                text = item["text"]
                full_content_parts.append(f"--- Slide {p_num} ---\n{text}\n")
                
            full_content = "\n".join(full_content_parts)
            return len(items), full_content, items

        elif ext in [".txt", ".md"]:
            try:
                print(f"Reading text/markdown file: {filename}...")
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
                items = [{
                    "text": text.strip(),
                    "page_number": 1
                }]
                full_content = text
                return 1, full_content, items
            except Exception as e:
                print(f"Failed to parse text file {file_path}: {e}")
                return 0, "", []
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def query(self, query_text: str, subject: str = None, limit: int = 5):
        """
        Query system for educational context.
        Uses pure-Python keyword overlap + frequency relevance ranking to find the top pages.
        """
        pages = db.get_all_pages(subject=subject)
        if not pages:
            return []
            
        # Tokenize query, ignore small words unless query is very short
        query_words = [w.lower() for w in query_text.split() if len(w) > 2]
        if not query_words:
            query_words = [w.lower() for w in query_text.split()]
            
        # Stop words filter for basic search
        stopwords = {"what", "how", "why", "who", "when", "where", "which", "the", "and", "for", "with", "about"}
        query_words = [w for w in query_words if w not in stopwords]
        
        # If query becomes empty after stop words, revert to original
        if not query_words:
            query_words = [w.lower() for w in query_text.split()]

        scored_pages = []
        for p in pages:
            content = p["content"]
            content_lower = content.lower()
            
            # Simple keyword matching score
            score = 0.0
            for word in query_words:
                count = content_lower.count(word)
                if count > 0:
                    # Score is weighted by exact matches and frequency
                    score += 1.0 + (0.1 * count)
                    
            if score > 0:
                scored_pages.append((score, p))
                
        # Sort by score descending
        scored_pages.sort(key=lambda x: x[0], reverse=True)
        
        # Select best matches
        if not scored_pages:
            # If no matches, return the most recently uploaded pages up to limit
            results_pages = pages[-limit:]
        else:
            results_pages = [p for score, p in scored_pages[:limit]]
            
        formatted_results = []
        for p in results_pages:
            formatted_results.append({
                "content": p["content"],
                "metadata": {
                    "filename": p["filename"],
                    "subject": p["subject"],
                    "page": p["page_number"],
                    "source": f"{p['filename']} (Page/Slide {p['page_number']})"
                }
            })
        return formatted_results
